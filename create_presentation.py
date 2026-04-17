from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Цвета из твоей презентации
BG_DARK = RGBColor(15, 23, 42)      # #0f172a
BG_ELEVATED = RGBColor(10, 22, 40)  # #0a1628 (из скриншота)
TEXT_PRIMARY = RGBColor(241, 245, 249)  # #f1f5f9
TEXT_SECONDARY = RGBColor(203, 213, 225)  # #cbd5e1
PRIMARY_BLUE = RGBColor(0, 212, 255)  # #00d4ff (из скриншота)
ACCENT_ORANGE = RGBColor(255, 107, 90)  # #ff6b5a (из скриншота)
SUCCESS_GREEN = RGBColor(0, 217, 165)  # #00d9a5
BORDER_ORANGE = RGBColor(255, 140, 122)  # для рамок

# Создаём презентацию
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

def add_title_slide(prs, title, subtitle="", tagline=""):
    """Титульный слайд"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = BG_ELEVATED

    # Рамка (как на скриншоте)
    border = slide.shapes.add_shape(
        1,  # Rectangle
        Inches(0.3), Inches(0.3),
        Inches(9.4), Inches(6.9)
    )
    border.fill.background()
    border.line.color.rgb = ACCENT_ORANGE
    border.line.width = Pt(4)

    # Заголовок
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1.5))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(64)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_BLUE
    p.alignment = PP_ALIGN.CENTER

    # Подзаголовок
    if subtitle:
        subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(8), Inches(1.2))
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.word_wrap = True
        p = subtitle_frame.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(18)
        p.font.color.rgb = TEXT_SECONDARY
        p.alignment = PP_ALIGN.CENTER

    # Таглайн
    if tagline:
        tag_box = slide.shapes.add_textbox(Inches(1), Inches(6.5), Inches(8), Inches(0.8))
        tag_frame = tag_box.text_frame
        p = tag_frame.paragraphs[0]
        p.text = tagline
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(148, 163, 184)
        p.alignment = PP_ALIGN.CENTER

def add_content_slide(prs, title, content_items, accent_color=PRIMARY_BLUE):
    """Слайд с контентом"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = BG_DARK

    # Рамка
    border = slide.shapes.add_shape(
        1, Inches(0.3), Inches(0.3),
        Inches(9.4), Inches(6.9)
    )
    border.fill.background()
    border.line.color.rgb = PRIMARY_BLUE
    border.line.width = Pt(3)

    # Заголовок
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.6), Inches(8.4), Inches(0.8))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = accent_color

    # Контент
    content_box = slide.shapes.add_textbox(Inches(1.2), Inches(1.8), Inches(7.6), Inches(5))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True

    for i, item in enumerate(content_items):
        if i > 0:
            p = text_frame.add_paragraph()
        else:
            p = text_frame.paragraphs[0]

        # Проверяем, есть ли маркер
        if item.startswith("•"):
            p.text = item[1:].strip()
            p.level = 0
        elif item.startswith("  -"):
            p.text = item[3:].strip()
            p.level = 1
        else:
            p.text = item

        p.font.size = Pt(20)
        p.font.color.rgb = TEXT_SECONDARY
        p.space_before = Pt(8)
        p.space_after = Pt(8)

def add_two_column_slide(prs, title, left_content, right_content):
    """Слайд с двумя колонками"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = BG_DARK

    # Рамка
    border = slide.shapes.add_shape(
        1, Inches(0.3), Inches(0.3),
        Inches(9.4), Inches(6.9)
    )
    border.fill.background()
    border.line.color.rgb = PRIMARY_BLUE
    border.line.width = Pt(3)

    # Заголовок
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.6), Inches(8.4), Inches(0.8))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_BLUE

    # Левая колонка
    left_box = slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(4), Inches(5))
    left_frame = left_box.text_frame
    left_frame.word_wrap = True
    for i, item in enumerate(left_content):
        if i > 0:
            p = left_frame.add_paragraph()
        else:
            p = left_frame.paragraphs[0]
        p.text = item
        p.font.size = Pt(18)
        p.font.color.rgb = TEXT_SECONDARY
        p.space_before = Pt(6)

    # Правая колонка
    right_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.8), Inches(4), Inches(5))
    right_frame = right_box.text_frame
    right_frame.word_wrap = True
    for i, item in enumerate(right_content):
        if i > 0:
            p = right_frame.add_paragraph()
        else:
            p = right_frame.paragraphs[0]
        p.text = item
        p.font.size = Pt(18)
        p.font.color.rgb = TEXT_SECONDARY
        p.space_before = Pt(6)

# ===== СЛАЙД 1: ТИТУЛ =====
add_title_slide(prs,
    "10 Ideas Bot",
    "Как я автоматизировала креативный процесс",
    "Презентация и Автоматизация в n8n • Telegram"
)

# ===== СЛАЙД 2: ПРОБЛЕМА =====
add_content_slide(prs, "Проблема", [
    "• Нужна регулярная практика развития креативности",
    "• Упражнение '10 идей' — но упражнение не принималось само собой",
    "",
    "↓",
    "",
    "Решение",
    "• Бот, который каждый день присылает 10 актуальных идей для креативности",
    "• 10 актуальных идей для решения реальных проблем"
], ACCENT_ORANGE)

# ===== СЛАЙД 3: ПЕРВАЯ ПОПЫТКА — ПРОВАЛ =====
add_content_slide(prs, "Первая версия: 100% провал", [
    "Метрика: 100% заданий доходят до 5-й попытки",
    "",
    "Проблема: Similarity score застревает на 88%",
    "",
    "Причина: Критик слишком строг (пороги 70-80)",
    "",
    "Результат: Всегда срабатывает fallback с простым промптом"
], ACCENT_ORANGE)

# ===== СЛАЙД 4: РЕШЕНИЕ =====
add_two_column_slide(prs, "Решение: Генератор + Критик", [
    "Генератор (DeepSeek):",
    "• 20 социальных доменов",
    "• Социальная алхимия",
    "• Микро-власть",
    "• Информационное выживание",
    "• Знаниевая инженерия",
    "• ..."
], [
    "Критик (Claude Haiku 4.5):",
    "• 7 правил проверки",
    "• История 10 заданий",
    "• Адаптивные пороги:",
    "  90 → 92 → 95",
    "• Emergency PASS на попытке 4+"
])

# ===== СЛАЙД 5: ТЕХНОЛОГИЯ =====
add_content_slide(prs, "Ключ к качеству: Адаптивные пороги", [
    "Правила критика:",
    "✅ 5+ уникальных слов совпадают → REJECT",
    "✅ Домен повторяется 5+ раз подряд → REJECT",
    "✅ Два абсолютных ограничения → REJECT",
    "⚠️  Одно ограничение (гипербола) → PASS",
    "",
    "Адаптивность:",
    "• Попытка 0-1: порог similarity ≥ 90",
    "• Попытка 2-3: порог ≥ 92",
    "• Попытка 4+: AUTO-PASS (emergency блок)"
], SUCCESS_GREEN)

# ===== СЛАЙД 6: WORKFLOW =====
add_content_slide(prs, "Как это работает: n8n Workflow", [
    "1. Daily Trigger 9AM → Получить активных пользователей",
    "",
    "2. Route: Первая генерация или регенерация?",
    "   ├─ regeneration_count = 0 → AI Генератор",
    "   ├─ regeneration_count 1-4 → AI Регенератор",
    "   └─ regeneration_count ≥ 5 → Fallback",
    "",
    "3. AI: Критик → PASS или REJECT?",
    "   ├─ PASS → Отправка пользователю",
    "   └─ REJECT → Log + возврат на регенерацию",
    "",
    "4. Save to Airtable Ideas"
])

# ===== СЛАЙД 7: ПРИМЕРЫ =====
add_content_slide(prs, "Реальный лог отклонений", [
    "Попытка 1: (REJECT)",
    '"10 способов убедить команду работать в выходные без принуждения"',
    '→ Критик: "Домен \'Социальная алхимия\' 6/8. Similarity 88."',
    "",
    "Попытка 2: (REJECT)",
    '"10 способов организовать проект, когда участники саботируют сроки"',
    '→ Критик: "Тот же домен. Вариация на \'враждебная среда\'. Similarity 88."',
    "",
    "Попытка 4: (AUTO-PASS) ✅",
    '→ Критик: "Emergency PASS. Задание уникально."'
])

# ===== СЛАЙД 8: РЕЗУЛЬТАТЫ =====
add_two_column_slide(prs, "Результаты: От провала к системе", [
    "До исправлений:",
    "❌ 100% заданий → fallback",
    "❌ Similarity застревает на 88%",
    "❌ Критик убивает всё",
    "",
    "",
    "Главный инсайт:",
    '"ИИ нужны границы,',
    'но границы должны быть',
    'адаптивными"'
], [
    "После исправлений:",
    "✅ 60-70% проходят на попытках 1-3",
    "✅ Emergency PASS гарантирует результат",
    "✅ Система калибрует границы,",
    "   не блокирует гиперболу"
])

# Сохраняем
prs.save('10_ideas_bot_presentation.pptx')
print("✅ Презентация создана: 10_ideas_bot_presentation.pptx")
